#!/usr/bin/env python3


import cachetools.func
import io
import os
import subprocess
import tempfile
import time
import traceback
import xml.etree.ElementTree as ET
import zipfile
from typing import List, Union

import my_mistral
import PyPDF2
import pandas as pd
from bs4 import BeautifulSoup
from pptx import Presentation


import my_log
import my_gemini
import utils


pandoc_cmd = 'pandoc'
catdoc_cmd = 'catdoc'


def fb2_to_text(data: bytes, ext: str = '', lang: str = '') -> str:
    """convert from fb2 or epub (bytes) and other types of books file to string"""
    if isinstance(data, str):
        with open(data, 'rb') as f:
            data = f.read()

    ext = ext.lower()
    if ext.startswith('.'):
        ext = ext[1:]

    if ext == 'pages':
        return convert_pages_to_md(data)
    elif ext == 'numbers':
        return convert_numbers_to_csv(data)

    input_file = utils.get_tmp_fname() + '.' + ext

    with open(input_file, 'wb') as f:
        f.write(data)

    book_type = ext

    if 'epub' in book_type:
        proc = subprocess.run([pandoc_cmd, '+RTS', '-M256M', '-RTS', '-f', 'epub', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    elif 'pptx' in book_type:
        text = read_pptx(input_file)
        utils.remove_file(input_file)
        return text
    elif 'docx' in book_type or 'dotx' in book_type:
        proc = subprocess.run([pandoc_cmd, '-f', '+RTS', '-M256M', '-RTS', 'docx', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    elif 'html' in book_type:
        proc = subprocess.run([pandoc_cmd, '-f', '+RTS', '-M256M', '-RTS', 'html', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    elif 'odt' in book_type:
        proc = subprocess.run([pandoc_cmd, '-f', '+RTS', '-M256M', '-RTS', 'odt', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    elif 'rtf' in book_type:
        proc = subprocess.run([pandoc_cmd, '-f', '+RTS', '-M256M', '-RTS', 'rtf', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    elif 'doc' in book_type:
        try:
            proc = subprocess.run([catdoc_cmd, input_file], stdout=subprocess.PIPE)
        except FileNotFoundError as e:
            my_log.log2(f'my_pandoc:fb2_to_text: {e}')
            proc = None
    elif 'pdf' in book_type or 'djvu' in book_type:
        if 'djvu' in book_type:
            input_file = convert_djvu2pdf(input_file)

        pdf_reader = PyPDF2.PdfReader(input_file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text()

        if not text or len(text) < 100:
            with open(input_file, 'rb') as f:
                data = f.read()
            text_ = my_mistral.ocr_pdf(data, timeout=300)
            if text_ and len(text_) > len(text):
                text = text_

        utils.remove_file(input_file)
        return text
    elif book_type in ('xlsx', 'ods', 'xls', 'numbers', 'xltx', 'xltm'):
        xls = pd.ExcelFile(io.BytesIO(data))
        result = ''
        for sheet in xls.sheet_names:
            csv = xls.parse(sheet_name=sheet).to_csv(index=False)
            result += f'\n\n{sheet}\n\n{csv}\n\n'

        utils.remove_file(input_file)

        return result
    elif 'fb2' in book_type:
        proc = subprocess.run([pandoc_cmd, '+RTS', '-M256M', '-RTS', '-f', 'fb2', '-t', 'gfm', input_file], stdout=subprocess.PIPE)
    else:
        utils.remove_file(input_file)
        try:
            result = data.decode('utf-8').replace(u'\xa0', u' ')
            return result
        except Exception as error:
            my_log.log2(f'my_pandoc:fb2_to_text other type error {error}')
            return ''

    utils.remove_file(input_file)

    try:
        output = proc.stdout.decode('utf-8', errors='replace').strip()
    except AttributeError:
        output = ''


    # в doc файле нет текста, но может есть картинки?
    if 'doc' in book_type and len(output) < 100:
        # try to read images
        images = extract_images_from_doc(data)
        if images:
            result = ''
            for image in images:
                text = my_gemini.ocr_page(image)
                if 'EMPTY' in text and len(text) < 20:
                    text = ''
                result += text + '\n\n'
            if len(result) > len(output):
                output = result


    return output


def read_pptx(input_file: str) -> str:
    """read pptx file"""
    prs = Presentation(input_file)
    text = ''
    for _, slide in enumerate(prs.slides):
        for shape in slide.shapes: 
            if hasattr(shape, "text"): 
                text += shape.text + '\n'
    return text


def convert_djvu2pdf(input_file: str) -> str:
    '''convert djvu to pdf and delete source file, return new file name'''
    output_file = input_file + '.pdf'
    subprocess.run(['ddjvu', '-format=pdf', input_file, output_file], check=True)
    utils.remove_file(input_file)
    return output_file


@cachetools.func.ttl_cache(maxsize=10, ttl=5 * 60)
def convert_markdown_to_document(text: str, output_format: str) -> bytes | str:
    """
    Converts Markdown-formatted text to a specified document format (DOCX or PDF).
    - For DOCX, uses Pandoc directly (md -> docx).
    - For PDF, uses a two-step process: Pandoc (md -> docx) and then LibreOffice (docx -> pdf).

    Assumes Pandoc and LibreOffice are installed and accessible in the system's PATH.

    Args:
        text (str): The input text in Markdown format, with LaTeX math enclosed in $..$ or $$..$$.
        output_format (str): The desired output format ('docx' or 'pdf').

    Returns:
        bytes: The binary content of the generated document file, or b'' on failure.
    """

    temp_input_md_file = None
    temp_output_file = None
    temp_intermediate_docx_file = None
    temp_pdf_out_dir = None

    try:
        with tempfile.NamedTemporaryFile(mode='w+', delete=False, suffix='.md', encoding='utf-8') as f_md:
            temp_input_md_file = f_md.name
            f_md.write(text)

        if output_format == 'docx':
            temp_output_file = utils.get_tmp_fname() + '.docx'
            command = [
                'pandoc', '+RTS', '-M256M', '-RTS',
                '-f', 'markdown+hard_line_breaks', '-t', 'docx',
                '--highlight-style=pygments',
                '-o', temp_output_file, temp_input_md_file
            ]
            result = subprocess.run(command, capture_output=True, text=True, check=False)
            if result.returncode != 0:
                my_log.log2(f"my_pandoc:convert_markdown_to_document: Pandoc conversion to DOCX failed:\n{result.stderr}")
                return b''

        elif output_format == 'pdf':
            # Этап 1: Markdown -> DOCX с помощью Pandoc
            temp_intermediate_docx_file = utils.get_tmp_fname() + '.docx'
            pandoc_command = [
                'pandoc', '+RTS', '-M256M', '-RTS',
                '-f', 'markdown+hard_line_breaks', '-t', 'docx',
                '--highlight-style=pygments',
                '-o', temp_intermediate_docx_file, temp_input_md_file
            ]
            pandoc_result = subprocess.run(pandoc_command, capture_output=True, text=True, check=False)
            if pandoc_result.returncode != 0:
                my_log.log2(f"my_pandoc:convert_markdown_to_document: PDF generation failed at Pandoc stage (md->docx):\n{pandoc_result.stderr}")
                return b''

            # Этап 2: DOCX -> PDF с помощью LibreOffice
            if 'windows' in utils.platform().lower():
                exe = 'soffice.exe'
            else:
                exe = 'libreoffice'

            temp_pdf_out_dir = tempfile.mkdtemp()
            libreoffice_command = [
                exe,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_pdf_out_dir,
                temp_intermediate_docx_file
            ]
            result = subprocess.run(libreoffice_command, capture_output=True, text=True, check=False)

            if result.returncode != 0:
                my_log.log2(f"my_pandoc:convert_markdown_to_document: PDF generation failed at LibreOffice stage (docx->pdf):\n{result.stderr}")
                return b''

            # ВАЖНО: Пауза для Windows, чтобы система успела записать файл на диск
            if 'windows' in utils.platform().lower():
                time.sleep(5)

            pdf_basename = os.path.splitext(os.path.basename(temp_intermediate_docx_file))[0] + '.pdf'
            temp_output_file = os.path.join(temp_pdf_out_dir, pdf_basename)

            # ВАЖНО: Явная проверка, что LibreOffice действительно создал файл
            if not os.path.exists(temp_output_file):
                my_log.log2(f"my_pandoc:convert_markdown_to_document: LibreOffice did not create the expected PDF file at {temp_output_file}. Stderr: {result.stderr}")
                return b""
        else:
            return b''

        # Читаем финальный файл, если он существует
        if temp_output_file and os.path.exists(temp_output_file):
            with open(temp_output_file, 'rb') as f_doc:
                document_data = f_doc.read()
            return document_data
        else:
            my_log.log2(f'my_pandoc:convert_markdown_to_document: Output file was not created or found for format {output_format}')
            return b''

    except FileNotFoundError as e:
        my_log.log2(f"my_pandoc:convert_markdown_to_document: Required tool not found (Pandoc or LibreOffice). Please ensure they are installed in system's PATH: {e}")
        return b""
    except Exception as e:
        my_log.log2(f"my_pandoc:convert_markdown_to_document: An unexpected error occurred during conversion: {e}\n{traceback.format_exc()}")
        return b""
    finally:
        # Надежная очистка всех временных файлов и папок
        if temp_input_md_file:
            utils.remove_file(temp_input_md_file)
        if temp_intermediate_docx_file:
            utils.remove_file(temp_intermediate_docx_file)
        if temp_pdf_out_dir:
            utils.remove_dir(temp_pdf_out_dir)
        # temp_output_file для docx удаляется отдельно, если он создавался
        if output_format == 'docx' and temp_output_file:
            utils.remove_file(temp_output_file)


@cachetools.func.ttl_cache(maxsize=10, ttl=5 * 60)
def convert_html_to_docx(html: str) -> bytes:
    """
    Converts well-formed HTML (using correct, valid tags) to a DOCX document via LibreOffice (headless mode).
    Uses the exact working CLI command with explicit export filter:
    --convert-to "docx:Office Open XML Text"

    Assumes LibreOffice is installed. On Windows the function uses the full path:
    "C:\\Program Files\\LibreOffice\\program\\soffice.exe"
    On other platforms it uses "libreoffice" from PATH.

    Args:
        html (str): The input HTML string. Must be well-formed and use valid tags.

    Returns:
        bytes: The binary content of the generated DOCX, or b'' on failure.
    """
    temp_input_html_file = None
    temp_output_docx_file = None
    temp_output_dir = None

    try:
        # 1) Save HTML to a temporary file
        with tempfile.NamedTemporaryFile(mode='w+', delete=False, suffix='.html', encoding='utf-8') as f_html:
            temp_input_html_file = f_html.name
            f_html.write(html)

        # 2) Create a temporary output directory
        temp_output_dir = tempfile.mkdtemp()

        # 3) Build the LibreOffice command with explicit DOCX filter
        exe = r"C:\Program Files\LibreOffice\program\soffice.exe" if 'windows' in utils.platform().lower() else 'libreoffice'
        libreoffice_command = [
            exe,
            '--headless',
            '--convert-to', 'docx:Office Open XML Text',
            '--outdir', temp_output_dir,
            temp_input_html_file
        ]
        result = subprocess.run(libreoffice_command, capture_output=True, text=True, check=False)

        if result.returncode != 0:
            my_log.log2(f"my_pandoc:convert_html_to_docx: LibreOffice conversion failed (html->docx):\n{result.stderr}")
            return b''

        # 4) On Windows, allow a short delay to ensure the file is fully written
        if 'windows' in utils.platform().lower():
            time.sleep(3)

        # 5) Determine the path to the generated DOCX (same basename, .docx extension)
        basename = os.path.splitext(os.path.basename(temp_input_html_file))[0] + '.docx'
        temp_output_docx_file = os.path.join(temp_output_dir, basename)

        if not os.path.exists(temp_output_docx_file):
            my_log.log2(
                f"my_pandoc:convert_html_to_docx: LibreOffice did not create expected DOCX at {temp_output_docx_file}. "
                f"Stderr: {result.stderr}"
            )
            return b''

        # 6) Read and return the DOCX bytes
        with open(temp_output_docx_file, 'rb') as f_doc:
            return f_doc.read()

    except FileNotFoundError as e:
        my_log.log2(f"my_pandoc:convert_html_to_docx: LibreOffice not found or not accessible. Error: {e}")
        return b""
    except Exception as e:
        my_log.log2(f"my_pandoc:convert_html_to_docx: Unexpected error: {e}\n{traceback.format_exc()}")
        return b""
    finally:
        # Cleanup temp files/directories
        if temp_input_html_file:
            utils.remove_file(temp_input_html_file)
        if temp_output_docx_file:
            utils.remove_file(temp_output_docx_file)
        if temp_output_dir:
            utils.remove_dir(temp_output_dir)


def convert_html_to_pdf(text: str) -> bytes:
    """
    Converts well-formed HTML directly to PDF via LibreOffice (headless).
    Note: Direct HTML->PDF may render worse than HTML->DOCX->PDF.
    """
    temp_input_html_file = None
    temp_output_pdf_file = None
    temp_output_dir = None

    try:
        # 1) Save HTML to a temporary file
        with tempfile.NamedTemporaryFile(mode='w+', delete=False, suffix='.html', encoding='utf-8') as f_html:
            temp_input_html_file = f_html.name
            f_html.write(text)

        # 2) Create temporary output directory
        temp_output_dir = tempfile.mkdtemp()

        # 3) Run LibreOffice to convert HTML -> PDF
        exe = r"C:\Program Files\LibreOffice\program\soffice.exe" if 'windows' in utils.platform().lower() else 'libreoffice'
        libreoffice_command = [
            exe,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', temp_output_dir,
            temp_input_html_file
        ]
        result = subprocess.run(libreoffice_command, capture_output=True, text=True, check=False)
        if result.returncode != 0:
            my_log.log2(f"my_pandoc:convert_html_to_pdf (direct): LibreOffice failed (html->pdf):\n{result.stderr}")
            return b""

        # 4) Wait a bit on Windows
        if 'windows' in utils.platform().lower():
            time.sleep(3)

        # 5) Locate output PDF
        basename = os.path.splitext(os.path.basename(temp_input_html_file))[0] + '.pdf'
        temp_output_pdf_file = os.path.join(temp_output_dir, basename)
        if not os.path.exists(temp_output_pdf_file):
            my_log.log2(f"my_pandoc:convert_html_to_pdf (direct): Expected PDF not found at {temp_output_pdf_file}. Stderr: {result.stderr}")
            return b""

        # 6) Read and return bytes
        with open(temp_output_pdf_file, 'rb') as f_pdf:
            return f_pdf.read()

    except FileNotFoundError as e:
        my_log.log2(f"my_pandoc:convert_html_to_pdf (direct): LibreOffice not found. Error: {e}")
        return b""
    except Exception as e:
        my_log.log2(f"my_pandoc:convert_html_to_pdf (direct): Unexpected error: {e}\n{traceback.format_exc()}")
        return b""
    finally:
        if temp_input_html_file:
            utils.remove_file(temp_input_html_file)
        if temp_output_pdf_file:
            utils.remove_file(temp_output_pdf_file)
        if temp_output_dir:
            utils.remove_dir(temp_output_dir)


def convert_file_to_html(data: bytes, filename: str) -> str:
    """
    Convert any supported file to HTML, determining the input format from the filename extension.

    Args:
        data: The file content as bytes.
        filename: The name of the file, used to determine the input format.

    Returns:
        The converted content in HTML format as a string.
    """
    output_file: str = utils.get_tmp_fname() + '.html'  # Generate a temporary file name for the output
    result: str = ''
    _, file_extension = os.path.splitext(filename)
    input_format: str = file_extension[1:].lower()  # Remove the leading dot and convert to lowercase

    if input_format == 'txt':
        # autodetect codepage and convert to utf8

        data = utils.extract_text_from_bytes(data)
        if not data:
            my_log.log2(f'my_pandoc:convert_file_to_html: convert_file_to_html: no data or unknown codepage {filename}')
            return ''

    # Mapping of file extensions to pandoc input formats
    format_mapping: dict[str, str] = {
        'bib': 'biblatex',
        'bibtex': 'bibtex',
        'bits': 'bits',
        'commonmark': 'commonmark',
        'cm': 'commonmark_x',  # Assuming .cm is a common extension for CommonMark
        'creole': 'creole',
        'csljson': 'csljson',
        'csv': 'csv',
        'djot': 'djot',
        'docbook': 'docbook',
        'docx': 'docx',
        'dokuwiki': 'dokuwiki',
        'endnote': 'endnotexml', # Assuming .endnote is a possible extension
        'epub': 'epub',
        'fb2': 'fb2',
        'gfm': 'gfm',
        'haddock': 'haddock',
        'html': 'html',
        'htm': 'html',
        'xhtml': 'html',
        'ipynb': 'ipynb',
        'jats': 'jats',
        'jira': 'jira',
        'json': 'json',
        'tex': 'latex',
        'latex': 'latex',
        'man': 'man',
        'md': 'markdown',
        'markdown': 'markdown',
        'markdown_github': 'markdown_github',
        'mmd': 'markdown_mmd', # Assuming .mmd is a common extension for MultiMarkdown
        'markdown_phpextra': 'markdown_phpextra',
        'markdown_strict': 'markdown_strict',
        'mediawiki': 'mediawiki',
        'muse': 'muse',
        'native': 'native',
        'odt': 'odt',
        'opml': 'opml',
        'org': 'org',
        'ris': 'ris',
        'rst': 'rst',
        'rtf': 'rtf',
        't2t': 't2t',
        'textile': 'textile',
        'txt': 'commonmark',
        'tikiwiki': 'tikiwiki',
        'tsv': 'tsv',
        'twiki': 'twiki',
        'typst': 'typst',
        'vimwiki': 'vimwiki',
    }

    pandoc_format: str | None = format_mapping.get(input_format)

    if not pandoc_format:
        my_log.log2(f'my_pandoc:convert_file_to_html: Unsupported file extension - {input_format}, defaulting to commonmark.')
        pandoc_format = 'commonmark' # Вот здесь мы устанавливаем формат на 'commonmark'

    try:
        # Execute the pandoc command to convert the file
        process = subprocess.run(
            ['pandoc', '+RTS', '-M256M', '-RTS', '-f', pandoc_format, '-t', 'html', '-o', output_file, '-'],
            input=data,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,  # Capture standard error for potential issues
            check=True # Raise an exception for non-zero exit codes
        )
        # Check for errors in pandoc execution
        if process.stderr:
            my_log.log2(f'my_pandoc:convert_file_to_html: Pandoc error - {process.stderr.decode()}')
            return "" # Or handle the error as needed

        with open(output_file, 'r', encoding='utf-8') as f:
            result = f.read()
    except FileNotFoundError:
        my_log.log2('my_pandoc:convert_file_to_html: Pandoc not found. Ensure it is installed and in your PATH.')
    except subprocess.CalledProcessError as error:
        my_log.log2(f'my_pandoc:convert_file_to_html: Pandoc conversion failed - {error}')
    except Exception as error:
        my_log.log2(f'my_pandoc:convert_file_to_html: An unexpected error occurred - {error}')
    finally:
        utils.remove_file(output_file)  # Clean up the temporary file
    return result


def ensure_utf8_meta(html_content: str) -> str:
    """
    Ensures the HTML content has a UTF-8 charset meta tag using BeautifulSoup,
    replacing any existing charset declaration.

    Args:
        html_content: The HTML content as a string.

    Returns:
        The HTML content with a UTF-8 charset meta tag.
    """
    soup = BeautifulSoup(html_content, 'html.parser')
    meta_charset = soup.find('meta', attrs={'charset': True})

    if meta_charset:
        meta_charset['charset'] = 'utf-8'
    else:
        meta_charset = soup.new_tag('meta', charset='utf-8')
        head_tag = soup.find('head')
        if head_tag:
            head_tag.insert(0, meta_charset)
        else:
            soup.insert(0, meta_charset) # insert to the beginning if no <head> tag

    return str(soup)


def convert_html_to_bytes(html_data: str, output_filename: str) -> bytes:
    """
    Convert HTML content to bytes of a specified format, determining the output
    format from the output filename extension, using a temporary file.

    Args:
        html_data: The HTML content as a string.
        output_filename: The name of the output file, used to determine the output format.

    Returns:
        The converted content in bytes.
    """

    # Гарантируем, что в HTML задана кодировка UTF-8
    html_data = ensure_utf8_meta(html_data)

    _, file_extension = os.path.splitext(output_filename)
    output_format: str = file_extension[1:].lower()  # Remove the leading dot and convert to lowercase

    # Mapping of file extensions to pandoc output formats
    format_mapping_out: dict[str, str] = {
        'adoc': 'asciidoc',
        'asciidoc': 'asciidoc',
        'beamer': 'beamer',
        'bib': 'biblatex',
        'biblatex': 'biblatex',
        'bibtex': 'bibtex',
        'csljson': 'csljson',
        'csv': 'csv',
        'context': 'context',
        'djot': 'djot',
        'docbook': 'docbook',
        'docbook4': 'docbook4',
        'docbook5': 'docbook5',
        'docx': 'docx',
        'dokuwiki': 'dokuwiki',
        'dzslides': 'dzslides',
        'epub': 'epub',
        'epub2': 'epub2',
        'epub3': 'epub3',
        'fb2': 'fb2',
        'gfm': 'gfm',
        'haddock': 'haddock',
        'html': 'html',
        'htm': 'html',
        'html4': 'html4',
        'html5': 'html5',
        'icml': 'icml',
        'ipynb': 'ipynb',
        'jats': 'jats',
        'jats_archiving': 'jats_archiving',
        'jats_articleauthoring': 'jats_articleauthoring',
        'jats_publishing': 'jats_publishing',
        'jira': 'jira',
        'json': 'json',
        'tex': 'latex',
        'latex': 'latex',
        'man': 'man',
        'md': 'markdown',
        'markdown': 'markdown',
        'markdown_github': 'markdown_github',
        'markdown_mmd': 'markdown_mmd',
        'markdown_phpextra': 'markdown_phpextra',
        'markdown_strict': 'markdown_strict',
        'markua': 'markua',
        'mediawiki': 'mediawiki',
        'ms': 'ms',
        'muse': 'muse',
        'native': 'native',
        'odt': 'odt',
        'opendocument': 'opendocument',
        'opml': 'opml',
        'org': 'org',
        'pdf': 'pdf',
        'pptx': 'pptx',
        'plain': 'plain',
        'revealjs': 'revealjs',
        'rst': 'rst',
        'rtf': 'rtf',
        's5': 's5',
        'slideous': 'slideous',
        'slidy': 'slidy',
        'tei': 'tei',
        'txt': 'plain',
        'texinfo': 'texinfo',
        'textile': 'textile',
        'typst': 'typst',
        'xwiki': 'xwiki',
        'zimwiki': 'zimwiki',
    }

    pandoc_format_out: str | None = format_mapping_out.get(output_format)

    if not pandoc_format_out:
        my_log.log2(f'my_pandoc:convert_html_to_bytes:1: Unsupported output file extension - {output_format}, defaulting to plain text.')
        pandoc_format_out = 'plain' # Вот здесь мы устанавливаем формат на 'plain' для неподдерживаемых расширений

    temp_output_file: str = utils.get_tmp_fname()  # Generate a temporary file name
    try:
        # Execute the pandoc command to convert the HTML to a temporary file
        process = subprocess.run(
            ['pandoc', '+RTS', '-M256M', '-RTS', '-f', 'html', '-t', pandoc_format_out, '-o', temp_output_file, '-'],
            input=html_data.encode('utf-8', 'replace'),  # Encode HTML string to bytes
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,  # Capture standard error for potential issues
            check=True  # Raise an exception for non-zero exit codes
        )
        # Check for errors in pandoc execution
        if process.stderr:
            my_log.log2(f'my_pandoc:convert_html_to_bytes:2: Pandoc error - {process.stderr.decode()}')
            # return b""

        with open(temp_output_file, 'rb') as f:
            return f.read()

    except FileNotFoundError:
        my_log.log2('my_pandoc:convert_html_to_bytes:3: Pandoc not found. Ensure it is installed and in your PATH.')
        return b""
    except subprocess.CalledProcessError as error:
        my_log.log2(f'my_pandoc:convert_html_to_bytes:4: Pandoc conversion failed - {error}')
        return b""
    except Exception as error:
        my_log.log2(f'my_pandoc:convert_html_to_bytes:5: An unexpected error occurred - {error}')
        return b""
    finally:
        utils.remove_file(temp_output_file)  # Clean up the temporary file


def convert_html_to_plain(html_data: str) -> str:
    """
    Converts HTML data to plain text using pandoc.

    Args:
        html_data: The HTML string to convert.

    Returns:
        The plain text representation of the HTML data.
    """
    result: str = html_data

    try:
        src_file: str = utils.get_tmp_fname() + '.html'
        dst_file: str = utils.get_tmp_fname() + '.txt'

        with open(src_file, 'w', encoding='utf-8', errors='replace') as f:
            f.write(html_data)

        process = subprocess.run(
            ['pandoc', '+RTS', '-M256M', '-RTS', '-f', 'html', '-t', 'plain', '-o', dst_file, src_file],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE)

        if process.returncode != 0:
            error_message = f"Pandoc failed with return code: {process.returncode}\nstderr: {process.stderr.decode('utf-8', errors='replace')}"
            my_log.log2(f'my_pandoc:convert_html_to_plain: {error_message}')
        else:
            with open(dst_file, 'r', encoding='utf-8', errors='replace') as f:
                result = f.read()

    except Exception as error:
        traceback_error = traceback.format_exc()
        my_log.log2(f'my_pandoc:convert_html_to_plain: {error}\n{traceback_error}') # Log other exceptions

    utils.remove_file(src_file)
    utils.remove_file(dst_file)

    return result


def convert_pages_to_md(pages_data: bytes) -> str:
    """
    Converts a .pages file (provided as bytes) to Markdown format (string).
    Uses LibreOffice to convert .pages to .odt, then Pandoc to convert .odt to .md.

    Args:
        pages_data (bytes): The content of the .pages file as bytes.

    Returns:
        str: The converted content in Markdown format, or an empty string if conversion fails.
    """
    temp_pages_file = None
    temp_odt_file = None
    markdown_content = ""

    try:
        if isinstance(pages_data, str):
            with open(pages_data, 'rb') as f:
                pages_data = f.read()

        # 1. Create a temporary .pages file from the input bytes
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pages") as f_pages:
            temp_pages_file = f_pages.name
            f_pages.write(pages_data)

        # Determine the expected output .odt file path
        # LibreOffice will create a file with the same base name as the input, but with .odt extension,
        # in the directory specified by --outdir.
        temp_dir = os.path.dirname(temp_pages_file)
        pages_base_name = os.path.basename(temp_pages_file)
        # Ensure correct .odt extension, handling cases where suffix might not be exact .pages
        odt_base_name = os.path.splitext(pages_base_name)[0] + ".odt"
        temp_odt_file = os.path.join(temp_dir, odt_base_name)

        # 2. Convert .pages to .odt using LibreOffice
        if 'windows' in utils.platform().lower():
            exe = 'soffice.exe'
        else:
            exe = 'libreoffice'

        libreoffice_cmd = [
            exe,
            "--headless",  # Run without GUI
            "--convert-to", "odt",
            "--outdir", temp_dir,  # Specify output directory
            temp_pages_file  # Input file
        ]

        libreoffice_process = subprocess.run(
            libreoffice_cmd,
            capture_output=True,  # Capture stdout and stderr
            text=True,            # Decode stdout/stderr as text
            check=False           # Don't raise CalledProcessError automatically
        )

        if libreoffice_process.returncode != 0:
            my_log.log2(f"my_pandoc:convert_pages_to_md: LibreOffice conversion to ODT failed with error code {libreoffice_process.returncode}:\n{libreoffice_process.stderr}")
            return ""

        if 'windows' in utils.platform().lower():
             time.sleep(5)

        # Check if the .odt file was actually created by LibreOffice
        if not os.path.exists(temp_odt_file):
             my_log.log2(f"my_pandoc:convert_pages_to_md: LibreOffice did not create the expected ODT file at {temp_odt_file}. Stderr: {libreoffice_process.stderr}")
             return ""

        # 3. Convert .odt to Markdown using Pandoc
        pandoc_cmd = [
            "pandoc",
            "+RTS", "-M256M", "-RTS",  # Memory limit for pandoc
            "-f", "odt",               # Input format: ODT
            "-t", "gfm",               # Output format: GitHub Flavored Markdown
            temp_odt_file              # Input file
        ]

        pandoc_process = subprocess.run(
            pandoc_cmd,
            capture_output=True,
            text=True,
            check=False
        )

        if pandoc_process.returncode != 0:
            my_log.log2(f"my_pandoc:convert_pages_to_md: Pandoc conversion to Markdown failed with error code {pandoc_process.returncode}:\n{pandoc_process.stderr}")
            return ""

        markdown_content = pandoc_process.stdout

    except FileNotFoundError as e:
        my_log.log2(f"my_pandoc:convert_pages_to_md: Required tool not found (LibreOffice or Pandoc). Please ensure they are installed and in your system's PATH: {e}")
        markdown_content = ""
    except Exception as e:
        my_log.log2(f"my_pandoc:convert_pages_to_md: An unexpected error occurred during conversion: {e}\n{traceback.format_exc()}")
        markdown_content = ""
    finally:
        # Clean up temporary files
        if temp_pages_file and os.path.exists(temp_pages_file):
            utils.remove_file(temp_pages_file)
        if temp_odt_file and os.path.exists(temp_odt_file):
            utils.remove_file(temp_odt_file)

    return markdown_content


def convert_numbers_to_csv(numbers_data: bytes) -> str:
    """
    Converts a .numbers file (provided as bytes) to CSV format (string), handling multiple sheets.
    Uses LibreOffice to convert .numbers to .ods, then pandas to read all sheets from .ods
    and convert them to CSV strings.

    Args:
        numbers_data (bytes): The content of the .numbers file as bytes.

    Returns:
        str: The concatenated CSV content of all sheets, separated by sheet names,
             or an empty string if conversion fails.
    """
    temp_numbers_file = None
    temp_ods_file = None # Изменим на .ods
    csv_content_all_sheets = [] # Для сбора CSV со всех листов

    try:
        if isinstance(numbers_data, str):
            # If a string (likely a file path) is passed, read its content
            with open(numbers_data, 'rb') as f:
                numbers_data = f.read()

        # 1. Create a temporary .numbers file from the input bytes
        with tempfile.NamedTemporaryFile(delete=False, suffix=".numbers") as f_numbers:
            temp_numbers_file = f_numbers.name
            f_numbers.write(numbers_data)

        # Determine the expected output .ods file path
        temp_dir = os.path.dirname(temp_numbers_file)
        numbers_base_name = os.path.basename(temp_numbers_file)
        # LibreOffice names the output file based on the input file's base name
        ods_base_name = os.path.splitext(numbers_base_name)[0] + ".ods"
        temp_ods_file = os.path.join(temp_dir, ods_base_name)

        # Determine the LibreOffice executable name based on the OS
        if 'windows' in utils.platform().lower():
            exe = 'soffice.exe'
        else:
            exe = 'libreoffice'

        # 2. Convert .numbers to .ods using LibreOffice
        libreoffice_cmd = [
            exe,
            "--headless",           # Run without GUI
            "--convert-to", "ods",  # Convert to ODS (OpenDocument Spreadsheet)
            "--outdir", temp_dir,   # Specify output directory
            temp_numbers_file       # Input file
        ]

        libreoffice_process = subprocess.run(
            libreoffice_cmd,
            capture_output=True,  # Capture stdout and stderr
            text=True,            # Decode stdout/stderr as text
            check=False           # Don't raise CalledProcessError automatically
        )

        if libreoffice_process.returncode != 0:
            my_log.log2(f"my_pandoc:convert_numbers_to_csv: LibreOffice conversion to ODS failed with error code {libreoffice_process.returncode}:\n{libreoffice_process.stderr}")
            return ""

        # Add a sleep for Windows, as LibreOffice might be delayed in writing the file
        if 'windows' in utils.platform().lower():
             time.sleep(5) # Give it 5 seconds to write the file

        # Check if the .ods file was actually created by LibreOffice
        if not os.path.exists(temp_ods_file):
             my_log.log2(f"my_pandoc:convert_numbers_to_csv: LibreOffice did not create the expected ODS file at {temp_ods_file}. Stderr: {libreoffice_process.stderr}")
             return ""

        # 3. Read the .ods file using pandas and convert each sheet to CSV
        try:
            # sheet_name=None reads all sheets into a dictionary of DataFrames
            all_sheets = pd.read_excel(temp_ods_file, engine='odf', sheet_name=None)

            for sheet_name, df in all_sheets.items():
                csv_buffer = io.StringIO()
                df.to_csv(csv_buffer, index=False, encoding='utf-8')
                # Добавляем заголовок для каждого листа и его содержимое
                csv_content_all_sheets.append(f"### Sheet: {sheet_name}\n")
                csv_content_all_sheets.append(csv_buffer.getvalue())
                csv_content_all_sheets.append("\n\n") # Добавляем разделитель между листами

        except Exception as e:
            my_log.log2(f"my_pandoc:convert_numbers_to_csv: Error reading ODS with pandas or converting to CSV: {e}\n{traceback.format_exc()}")
            return ""

    except FileNotFoundError as e:
        my_log.log2(f"my_pandoc:convert_numbers_to_csv: Required tool (LibreOffice) not found. Please ensure it is installed and in your system's PATH: {e}")
        csv_content_all_sheets = []
    except Exception as e:
        my_log.log2(f"my_pandoc:convert_numbers_to_csv: An unexpected error occurred during conversion: {e}\n{traceback.format_exc()}")
        csv_content_all_sheets = []
    finally:
        # Clean up temporary files
        if temp_numbers_file and os.path.exists(temp_numbers_file):
            utils.remove_file(temp_numbers_file)
        if temp_ods_file and os.path.exists(temp_ods_file):
            utils.remove_file(temp_ods_file)

    return "".join(csv_content_all_sheets)


def extract_images_from_doc(doc_data: Union[bytes, str]) -> List[bytes]:
    """
    Extracts images from a .doc file in their original order.

    The process involves:
    1. Converting the .doc file to .docx using LibreOffice.
    2. Reading the resulting .docx file (which is a zip archive).
    3. Parsing word/_rels/document.xml.rels to map relationship IDs to image files.
    4. Parsing word/document.xml to find image references in their sequential order.
    5. Extracting the image bytes from the archive based on the ordered references.

    Args:
        doc_data: The content of the .doc file as bytes or a path to the file as a string.

    Returns:
        A list of bytes, where each item is the binary content of an image,
        sorted in the order of their appearance in the document. Returns an empty
        list if conversion or extraction fails.
    """
    temp_dir = tempfile.mkdtemp()
    temp_doc_path = os.path.join(temp_dir, "input.doc")
    images_bytes: List[bytes] = []

    try:
        # Step 1: Write input data to a temporary .doc file
        if isinstance(doc_data, str):
            with open(doc_data, 'rb') as f:
                doc_content = f.read()
        else:
            doc_content = doc_data

        with open(temp_doc_path, 'wb') as f:
            f.write(doc_content)

        # Step 2: Convert .doc to .docx using LibreOffice
        exe = 'soffice.exe' if 'windows' in utils.platform().lower() else 'libreoffice'
        libreoffice_cmd = [
            exe,
            "--headless",
            "--convert-to", "docx",
            "--outdir", temp_dir,
            temp_doc_path
        ]

        result = subprocess.run(
            libreoffice_cmd, capture_output=True, text=True, check=False
        )

        if result.returncode != 0:
            my_log.log2(f"my_pandoc:extract_images_from_doc: LibreOffice failed: {result.stderr}")
            return []

        temp_docx_path = os.path.join(temp_dir, "input.docx")
        if not os.path.exists(temp_docx_path):
            my_log.log2("my_pandoc:extract_images_from_doc: DOCX output file not found after conversion.")
            return []

        # Step 3: Open the .docx as a zip archive and parse XMLs to get ordered images
        with zipfile.ZipFile(temp_docx_path, 'r') as docx_zip:
            # XML Namespace mapping
            ns = {
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
                'doc': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'draw': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }

            # Map Relationship IDs (rId) to image file paths (e.g., 'media/image1.png')
            id_to_target = {}
            try:
                rels_xml_content = docx_zip.read('word/_rels/document.xml.rels')
                rels_root = ET.fromstring(rels_xml_content)
                for rel in rels_root.findall('rel:Relationship', ns):
                    r_id = rel.get('Id')
                    target = rel.get('Target')
                    if r_id and target:
                        id_to_target[r_id] = target
            except KeyError:
                my_log.log2("my_pandoc:extract_images_from_doc: word/_rels/document.xml.rels not found.")
                return []

            # Find all image references ('r:embed' attributes) in the main document, preserving order
            ordered_rids = []
            try:
                doc_xml_content = docx_zip.read('word/document.xml')
                doc_root = ET.fromstring(doc_xml_content)
                # Find all 'blip' elements which contain the image embed reference
                for blip in doc_root.findall('.//draw:blip', ns):
                    embed_rid = blip.get(f'{{{ns["r"]}}}embed')
                    if embed_rid:
                        ordered_rids.append(embed_rid)
            except KeyError:
                my_log.log2("my_pandoc:extract_images_from_doc: word/document.xml not found.")
                return []

            # Extract image bytes from the archive using the ordered list of rIds
            for rid in ordered_rids:
                image_target = id_to_target.get(rid)
                if image_target and image_target.startswith('media/'):
                    image_path_in_zip = f'word/{image_target}'
                    try:
                        image_data = docx_zip.read(image_path_in_zip)
                        images_bytes.append(image_data)
                    except KeyError:
                        my_log.log2(f"my_pandoc:extract_images_from_doc: Image file not found in zip: {image_path_in_zip}")

    except Exception as e:
        my_log.log2(f"my_pandoc:extract_images_from_doc: An unexpected error occurred: {e}")
    finally:
        # Step 4: Clean up temporary directory and all its contents
        utils.remove_dir(temp_dir)

    return images_bytes


if __name__ == '__main__':
    pass
    # result = fb2_to_text(open('c:/Users/user/Downloads/1.xlsx', 'rb').read(), '.xlsx')
    # result = fb2_to_text(open('1.pdf', 'rb').read(), '.pdf')
    # print(result)
    # print(convert_djvu2pdf('/home/ubuntu/tmp/2.djvu'))

    # with open('c:/Users/user/Downloads/2.docx', 'rb') as f:
    #     html = convert_file_to_html(f.read(), '2.docx')
    #     with open('c:/Users/user/Downloads/2.html', 'w', encoding='utf-8') as f:
    #         f.write(html)

    # with open('c:/Users/user/Downloads/2.html', 'r', encoding='utf-8') as f:
    #     data = convert_html_to_bytes(f.read(), '3.docx')
    #     with open('c:/Users/user/Downloads/3.docx', 'wb') as f:
    #         f.write(data)

    # t = '''Громаш, словно обезумев, начал колотить кулаками по каменной стене, словно он пытался разбить её в щепки, словно он пытался вырваться из этой тюрьмы. &quot;Я найду тебя, Максимус, и я заставлю тебя страдать! - ревел он, и его голос разносился по всему залу, - я отплачу тебе за все мои страдания, и я не успокоюсь, пока не увижу твою смерть! <i>А было это так… Я буду мучить тебя так, как ты мучил моих братьев, я буду терзать тебя до тех пор, пока не вырву твою душу, и я буду наслаждаться твоими страданиями, словно это самый вкусный нектар. Я хочу видеть боль в твоих глазах, я хочу слышать твои крики, и это будет моей местью, это будет моим искуплением, и я не отступлю, пока я не добьюсь своего!</i>&quot; Он посмотрел на своих соплеменников, и его глаза горели адским огнем, словно он был одержим демоном. &quot;Мы покажем всему миру, что такое сила орков, и мы заставим всех бояться нас, и мы будем править этим миром, и все будут подчиняться нам, и все будут преклоняться перед нами!&quot; Орки, слыша его слова, начали вопить и стучать оружием, готовясь к битве, словно стая диких зверей, готовых разорвать свою жертву на куски. Максимус, слушая его крики, повернулся к своим товарищам и сказал: &quot;Мы все пали, но мы не должны сдаваться, мы не должны позволить этой ненависти сломить нас. Мы должны сражаться вместе, чтобы выжить, и мы должны помнить, что мы не одни.&quot; Его голос был тихим, но в нем чувствовалась решимость и готовность бороться до конца, словно он знал, что они смогут справиться со всеми испытаниями. &quot;Я не верю в шансы,&quot; - сказала Лираэль, и ее голос был полон цинизма, - &quot;но я готова сражаться, если это будет нужно для достижения моих целей, я готова убивать, если это поможет мне выжить, но я не верю в то, что у нас есть шанс, я не верю в то, что мы сможем изменить свою судьбу. <i>А было это так… Я всегда сражалась в одиночку, я не доверяла никому, и я всегда полагалась только на себя, но сейчас я чувствую, что что-то меняется, словно я становлюсь частью чего-то большего, и я не понимаю, что это значит. Я вспоминаю свои прошлые победы, как я достигала всего самостоятельно, но сейчас всё по-другому, и я не знаю, что будет дальше.</i>&quot; Она посмотрела на остальных, и в её глазах можно было заметить искорку надежды, словно она хотела верить, что у них есть шанс. Борд, вздохнув, проговорил: &quot;Я создал много мечей, но ни один из них не спас меня от предательства, и я не знаю, почему я должен сражаться, зачем я должен помогать другим, когда никто не помог мне, я просто устал, и я не знаю, что делать дальше. *А было это так… Я вспоминаю свою кузницу, я её так любил, я отдавал ей всю свою жизнь, но судьба забрала у меня всё, и я не знаю, почему это случилось со мной.'''
    # print(convert_html_to_plain(t))

    # r = convert_numbers_to_csv(r'c:\users\user\downloads\2.numbers')
    # print(r)

    # print(extract_images_from_doc(r'C:\Users\user\Downloads\Экзотическое_и_Громоздкое_оружие_1.doc'))
